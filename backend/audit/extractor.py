"""Slide extraction. Returns a list of Slide dicts."""

from dataclasses import dataclass, field
from typing import List, Optional
from pptx import Presentation


@dataclass
class Bullet:
    text: str
    level: int  # 0 = top level, 1+ = nested


@dataclass
class Slide:
    index: int               # 1-based slide number
    title: Optional[str]
    bullets: List[Bullet] = field(default_factory=list)
    body_text: str = ""      # all body text joined (excluding title)
    notes: str = ""
    word_count: int = 0

    @property
    def all_text(self) -> str:
        parts = []
        if self.title:
            parts.append(self.title)
        if self.body_text:
            parts.append(self.body_text)
        if self.notes:
            parts.append(self.notes)
        return "\n".join(parts)


def _shape_iter(slide):
    """Recursively yield shapes, descending into group shapes."""
    for shape in slide.shapes:
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            for inner in _walk_group(shape):
                yield inner
        else:
            yield shape


def _walk_group(group_shape):
    for shape in group_shape.shapes:
        if shape.shape_type == 6:
            for inner in _walk_group(shape):
                yield inner
        else:
            yield shape


def _title_text(slide) -> Optional[str]:
    try:
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            t = slide.shapes.title.text_frame.text.strip()
            return t or None
    except Exception:
        pass
    return None


def _notes_text(slide) -> str:
    try:
        if slide.has_notes_slide:
            return (slide.notes_slide.notes_text_frame.text or "").strip()
    except Exception:
        pass
    return ""


def extract(path: str) -> List[Slide]:
    prs = Presentation(path)
    title_shape_ids = set()
    out: List[Slide] = []

    for i, sl in enumerate(prs.slides, start=1):
        title = _title_text(sl)

        # Track the title placeholder so we don't re-grab it as body.
        title_id = None
        try:
            if sl.shapes.title is not None:
                title_id = sl.shapes.title.shape_id
        except Exception:
            pass

        bullets: List[Bullet] = []
        body_chunks: List[str] = []

        for shape in _shape_iter(sl):
            if title_id is not None and getattr(shape, "shape_id", None) == title_id:
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame
            for para in tf.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if not text:
                    # fall back to paragraph.text if runs were empty
                    text = (para.text or "").strip()
                if not text:
                    continue
                level = getattr(para, "level", 0) or 0
                bullets.append(Bullet(text=text, level=level))
                body_chunks.append(text)

        body_text = "\n".join(body_chunks)
        notes = _notes_text(sl)
        wc = len((title or "").split()) + len(body_text.split()) + len(notes.split())

        out.append(Slide(
            index=i,
            title=title,
            bullets=bullets,
            body_text=body_text,
            notes=notes,
            word_count=wc,
        ))

    return out
